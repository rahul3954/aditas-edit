import os
import re
import httpx
import json
from typing import List, Dict, Any
import pandas as pd
import pypdf
from io import BytesIO, StringIO
import re as stdlib_re
import smtplib
from email.message import EmailMessage
import imaplib
import email
import asyncio
import mimetypes

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Preformatted, Image
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_JUSTIFY, TA_LEFT, TA_CENTER
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
from app.db.database import SessionLocal
from app.db.models import SystemSetting
from app.db import models
from app.services.browser_automation import browser_automation

class MCPService:

    def __init__(self):
        # We now support real API keys if provided in .env
        self.tavily_api_key = os.getenv("TAVILY_API_KEY")
        self.serper_api_key = os.getenv("SERPER_API_KEY")
        self.smtp_email = os.getenv("SMTP_EMAIL")
        self.smtp_password = os.getenv("SMTP_PASSWORD")
        self.slack_webhook = os.getenv("SLACK_WEBHOOK_URL")
        
        self.tools = [
            # ... existing tools ...
            {
                "name": "google_search",
                "description": "Searches for high-level information, facts, and website URLs. Use this for general queries.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "query": {"type": "string", "description": "The search query."}
                    },
                    "required": ["query"]
                }
            },
            {
                "name": "browse_url",
                "description": "Visits a URL and extracts text content. For JS-heavy sites.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "url": {"type": "string", "description": "The URL to browse."}
                    },
                    "required": ["url"]
                }
            },
            {
                "name": "take_screenshot",
                "description": "Screenshots a website, saves to agent_files.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "url": {"type": "string", "description": "The URL to capture."},
                        "filename": {"type": "string", "description": "Filename for saving (e.g., 'amazon_home.png')."}
                    },
                    "required": ["url", "filename"]
                }
            },
            {
                "name": "write_file",
                "description": "Creates or overwrites a file with specific content in the 'agent_files' folder.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {"type": "string", "description": "Name of the file."},
                        "content": {"type": "string", "description": "Content to write."}
                    },
                    "required": ["filename", "content"]
                }
            },
            {
                "name": "analyze_data",
                "description": "Analyzes CSV/Excel data using Pandas. Returns stats and preview.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {"type": "string", "description": "Name of the file in agent_files (e.g., data.csv)."},
                        "query": {"type": "string", "description": "What to analyze (e.g., 'summary', 'total revenue')."}
                    },
                    "required": ["filename", "query"]
                }
            },
            {
                "name": "read_pdf",
                "description": "Extracts text from a PDF file located in the 'agent_files' folder.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {"type": "string", "description": "Name of the PDF file."}
                    },
                    "required": ["filename"]
                }
            },
            {
                "name": "draft_email",
                "description": "Drafts an email preview for user approval.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "recipient": {"type": "string", "description": "Email address of the recipient."},
                        "subject": {"type": "string", "description": "Subject line of the email."},
                        "body": {"type": "string", "description": "Body content of the email."},
                        "attachments": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "List of filenames in 'agent_files' to attach (e.g., ['report.pdf'])."
                        }
                    },
                    "required": ["recipient", "subject", "body"]
                }
            },
            {
                "name": "confirm_send_email",
                "description": "Sends the previously drafted email after user confirmation.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "confirmed": {"type": "boolean", "description": "User confirmation to send."}
                    },
                    "required": ["confirmed"]
                }
            },
            {
                "name": "schedule_task",
                "description": "Schedules a recurring task at an interval.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "task_description": {"type": "string", "description": "Description of the task (e.g., 'Check email')."},
                        "interval_seconds": {"type": "integer", "description": "Interval in seconds (e.g., 3600 for 1 hour)."}
                    },
                    "required": ["task_description", "interval_seconds"]
                }
            },
            {
                "name": "read_email",
                "description": "Fetches unread emails from the inbox via IMAP.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "limit": {"type": "integer", "description": "Number of emails to fetch (default 5)."}
                    },
                    "required": []
                }
            },
            {
                "name": "list_scheduled_tasks",
                "description": "Lists all currently active scheduled jobs/tasks.",
                "parameters": {
                    "type": "object",
                    "properties": {},
                    "required": []
                }
            },
            {
                "name": "cancel_task",
                "description": "Cancels a scheduled task by its Job ID.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "job_id": {"type": "string", "description": "The ID of the job to cancel."}
                    },
                    "required": ["job_id"]
                }
            },
            {
                "name": "create_pdf",
                "description": "Creates a PDF file with the given content.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {"type": "string", "description": "Name of the file (e.g., 'report.pdf')."},
                        "content": {"type": "string", "description": "Text content to write to the PDF."}
                    },
                    "required": ["filename", "content"]
                }
            },
            {
                "name": "create_docx",
                "description": "Creates a Word document (.docx) with the given content.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {"type": "string", "description": "Name of the file (e.g., 'report.docx')."},
                        "content": {"type": "string", "description": "Text content to write to the document."}
                    },
                    "required": ["filename", "content"]
                }
            },
            {
                "name": "create_ppt",
                "description": "Creates a PowerPoint presentation (.pptx) with slides.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {"type": "string", "description": "Name of the file (e.g., 'presentation.pptx')."},
                        "title": {"type": "string", "description": "Title of the presentation."},
                        "slides": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "title": {"type": "string"},
                                    "content": {"type": "array", "items": {"type": "string"}}
                                }
                            },
                            "description": "List of slides, each with a title and list of bullet points."
                        }
                    },
                    "required": ["filename", "title", "slides"]
                }
            },
            {
                "name": "create_excel",
                "description": "Creates an Excel spreadsheet (.xlsx).",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "filename": {"type": "string", "description": "Name of the file (e.g., 'data.xlsx')."},
                        "data": {
                            "type": "array",
                            "items": {"type": "object"},
                            "description": "List of dictionaries representing rows of data."
                        }
                    },
                    "required": ["filename", "data"]
                }
            },
            {
                "name": "generate_linkedin_post",
                "description": "Generates a professional LinkedIn-style post from a topic or description.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "topic": {"type": "string", "description": "Topic or description for the LinkedIn post."}
                    },
                    "required": ["topic"]
                }
            },
            {
                "name": "post_to_linkedin",
                "description": "Posts content to LinkedIn with optional images. Requires OAuth authentication.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "text": {"type": "string", "description": "The post text content."},
                        "image_filenames": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "Optional list of image filenames from agent_files to attach."
                        }
                    },
                    "required": ["text"]
                }
            },
            # === BROWSER AUTOMATION TOOLS (BrowserOS Architecture) ===
            # Pattern: open_browser → take_snapshot → interact by uid → take_snapshot
            {
                "name": "open_browser",
                "description": "Opens a browser and navigates to a URL. Automatically takes an initial snapshot. Browser stays open for further interactions.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "url": {"type": "string", "description": "The URL to navigate to."}
                    },
                    "required": ["url"]
                }
            },
            {
                "name": "take_snapshot",
                "description": "Takes a snapshot of the current page. Returns all interactive elements with unique IDs (uid). ALWAYS call this before interacting with elements. Use the uid values with click, fill, hover, etc.",
                "parameters": {
                    "type": "object",
                    "properties": {},
                    "required": []
                }
            },
            {
                "name": "click",
                "description": "Clicks an element by its uid from the latest snapshot. Automatically takes a new snapshot after clicking.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "uid": {"type": "string", "description": "The uid of the element from take_snapshot (e.g. '3_5')."},
                        "dbl_click": {"type": "boolean", "description": "Set to true for double-click. Default: false."}
                    },
                    "required": ["uid"]
                }
            },
            {
                "name": "hover",
                "description": "Hovers over an element by uid to trigger dropdowns/tooltips. Automatically takes a new snapshot showing revealed elements.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "uid": {"type": "string", "description": "The uid of the element from take_snapshot."}
                    },
                    "required": ["uid"]
                }
            },
            {
                "name": "fill",
                "description": "Fills a text input, textarea, or select element by uid. For select elements, pass the option label as value.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "uid": {"type": "string", "description": "The uid of the input element from take_snapshot."},
                        "value": {"type": "string", "description": "The value to fill in."}
                    },
                    "required": ["uid", "value"]
                }
            },
            {
                "name": "fill_form",
                "description": "Fills multiple form fields at once. More efficient than calling fill repeatedly.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "elements": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "uid": {"type": "string"},
                                    "value": {"type": "string"}
                                }
                            },
                            "description": "List of {uid, value} pairs to fill."
                        }
                    },
                    "required": ["elements"]
                }
            },
            {
                "name": "select_option",
                "description": "Selects an option from a dropdown by uid and option text.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "uid": {"type": "string", "description": "The uid of the select/dropdown element."},
                        "option_text": {"type": "string", "description": "The visible text of the option to select."}
                    },
                    "required": ["uid", "option_text"]
                }
            },
            {
                "name": "drag",
                "description": "Drags one element onto another by their uids.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "from_uid": {"type": "string", "description": "The uid of the element to drag."},
                        "to_uid": {"type": "string", "description": "The uid of the drop target."}
                    },
                    "required": ["from_uid", "to_uid"]
                }
            },
            {
                "name": "upload_file",
                "description": "Uploads a file to a file input element by uid.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "uid": {"type": "string", "description": "The uid of the file input element."},
                        "file_path": {"type": "string", "description": "Path to the file. Can be filename in agent_files or absolute path."}
                    },
                    "required": ["uid", "file_path"]
                }
            },
            {
                "name": "navigate_page",
                "description": "Navigates the current page to a URL. Automatically takes a snapshot.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "url": {"type": "string", "description": "The URL to navigate to."}
                    },
                    "required": ["url"]
                }
            },
            {
                "name": "navigate_history",
                "description": "Navigates back or forward in browser history.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "direction": {"type": "string", "description": "'back' or 'forward'."}
                    },
                    "required": ["direction"]
                }
            },
            {
                "name": "new_page",
                "description": "Opens a URL in a new tab. The new tab becomes active.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "url": {"type": "string", "description": "URL to open in the new tab."}
                    },
                    "required": ["url"]
                }
            },
            {
                "name": "list_pages",
                "description": "Lists all open pages/tabs with their indices.",
                "parameters": {
                    "type": "object",
                    "properties": {},
                    "required": []
                }
            },
            {
                "name": "select_page",
                "description": "Switches to a page/tab by index. Use list_pages to see available pages.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "page_idx": {"type": "integer", "description": "Page index (0-based)."}
                    },
                    "required": ["page_idx"]
                }
            },
            {
                "name": "close_page",
                "description": "Closes a page/tab by index. Cannot close the last page.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "page_idx": {"type": "integer", "description": "Page index to close."}
                    },
                    "required": ["page_idx"]
                }
            },
            {
                "name": "close_browser",
                "description": "Closes the browser session and cleans up.",
                "parameters": {
                    "type": "object",
                    "properties": {},
                    "required": []
                }
            },
            {
                "name": "take_screenshot",
                "description": "Takes a screenshot of the page or a specific element by uid.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "uid": {"type": "string", "description": "Optional uid of element to screenshot. If omitted, screenshots the full page."},
                        "full_page": {"type": "boolean", "description": "If true, captures the entire scrollable page. Default: false."}
                    },
                    "required": []
                }
            },
            {
                "name": "extract_text",
                "description": "Reads visible text from the current page.",
                "parameters": {
                    "type": "object",
                    "properties": {},
                    "required": []
                }
            },
            {
                "name": "extract_structured_data",
                "description": "Extracts tables/lists/headings/links as JSON. Types: 'auto', 'tables', 'lists', 'headings', 'links'.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "data_type": {"type": "string", "description": "Type: 'auto', 'tables', 'lists', 'headings', or 'links'. Default: 'auto'."}
                    },
                    "required": []
                }
            },
            {
                "name": "type_text",
                "description": "Types text char-by-char with human-like delays. Use for search boxes or React inputs.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "text": {"type": "string", "description": "The text to type."},
                        "uid": {"type": "string", "description": "Optional uid of element to type into. If omitted, types into focused element."}
                    },
                    "required": ["text"]
                }
            },
            {
                "name": "press_key",
                "description": "Presses a keyboard key or combination (e.g. Enter, Tab, Ctrl+A).",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "key": {"type": "string", "description": "Key: 'Enter', 'Tab', 'Escape', 'ArrowDown', etc."},
                        "modifiers": {"type": "string", "description": "Optional modifier: 'Control', 'Shift', 'Alt'."}
                    },
                    "required": ["key"]
                }
            },
            {
                "name": "scroll_page",
                "description": "Scrolls the page. Automatically takes a new snapshot.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "direction": {"type": "string", "description": "'up', 'down', 'top', or 'bottom'."}
                    },
                    "required": ["direction"]
                }
            },
            {
                "name": "scroll_to_element",
                "description": "Scrolls until a specific element (by uid) is visible.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "uid": {"type": "string", "description": "The uid of the element to scroll to."}
                    },
                    "required": ["uid"]
                }
            },
            {
                "name": "wait_for",
                "description": "Waits for specified text to appear on the page.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "text": {"type": "string", "description": "Text to wait for."},
                        "timeout": {"type": "integer", "description": "Max wait in ms (default 5000)."}
                    },
                    "required": ["text"]
                }
            },
            {
                "name": "wait_for_navigation",
                "description": "Waits for the page URL to change after a click or form submit.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "timeout": {"type": "integer", "description": "Max wait in ms (default 10000)."}
                    },
                    "required": []
                }
            },
            {
                "name": "execute_javascript",
                "description": "Executes JavaScript on the page and returns the result.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "code": {"type": "string", "description": "JavaScript code to execute."}
                    },
                    "required": ["code"]
                }
            },
            {
                "name": "handle_dialog",
                "description": "Handles browser alert/confirm/prompt dialogs.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "action": {"type": "string", "description": "'accept' or 'dismiss'. Default: 'accept'."},
                        "prompt_text": {"type": "string", "description": "Text for prompt dialogs (optional)."}
                    },
                    "required": []
                }
            },
            {
                "name": "switch_to_frame",
                "description": "Switches into an iframe by uid to interact with its content.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "uid": {"type": "string", "description": "The uid of the iframe element."}
                    },
                    "required": ["uid"]
                }
            },
            {
                "name": "switch_to_main",
                "description": "Switches back to the main page from an iframe.",
                "parameters": {
                    "type": "object",
                    "properties": {},
                    "required": []
                }
            },
            {
                "name": "get_page_info",
                "description": "Gets current page URL, title, tab count, and scroll position.",
                "parameters": {
                    "type": "object",
                    "properties": {},
                    "required": []
                }
            },
            {
                "name": "submit_form",
                "description": "Auto-finds and clicks a submit button. Prefer using click with a specific uid instead.",
                "parameters": {
                    "type": "object",
                    "properties": {},
                    "required": []
                }
            },
            # === YOUTUBE TRANSCRIPT TOOLS ===
            {
                "name": "youtube_transcript_search",
                "description": "Searches a YouTube video's transcript for a specific phrase and returns the timestamp where it appears. This is TOKEN-EFFICIENT: the search happens locally in Python, not via LLM. After getting the timestamp, navigate the browser to the video URL with ?t=seconds to play from that point.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "video_url": {"type": "string", "description": "YouTube video URL (e.g. https://www.youtube.com/watch?v=VIDEO_ID)"},
                        "search_phrase": {"type": "string", "description": "The phrase or sentence to search for in the transcript."}
                    },
                    "required": ["video_url", "search_phrase"]
                }
            },
            {
                "name": "get_youtube_transcript",
                "description": "Fetches the transcript of a YouTube video. Returns timestamped text. Use max_chars to limit output and save tokens. For searching specific phrases, use youtube_transcript_search instead.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "video_url": {"type": "string", "description": "YouTube video URL."},
                        "max_chars": {"type": "integer", "description": "Maximum characters to return (default 2000). Lower = fewer tokens."}
                    },
                    "required": ["video_url"]
                }
            }
        ]

    def get_tool_definitions(self) -> List[Dict[str, Any]]:
        return [
            {
                "function_declarations": [
                    {
                        "name": tool["name"],
                        "description": tool["description"],
                        "parameters": tool["parameters"]
                    } for tool in self.tools
                ]
            }
        ]

    async def execute_tool(self, name: str, arguments: Dict[str, Any]) -> str:
        """Executes the requested tool and returns the result."""
        try:
            if name == "google_search":
                return await self._real_search(arguments.get("query"))
            elif name == "browse_url":
                return await self._browse_url(arguments.get("url"))
            elif name == "take_screenshot":
                return await self._take_screenshot(arguments.get("url"), arguments.get("filename"))
            elif name == "write_file":
                return self._write_file(arguments.get("filename"), arguments.get("content"))
            elif name == "analyze_data":
                return self._analyze_data(arguments.get("filename"), arguments.get("query"))
            elif name == "read_pdf":
                return self._read_pdf(arguments.get("filename"))
            elif name == "draft_email":
                return self._draft_email(arguments.get("recipient"), arguments.get("subject"), arguments.get("body"), arguments.get("attachments"))
            elif name == "confirm_send_email":
                return self._confirm_send_email(arguments.get("confirmed"))
            elif name == "schedule_task":
                return self._schedule_task(arguments.get("task_description"), arguments.get("interval_seconds"))
            elif name == "read_email":
                return self._read_email(arguments.get("limit", 5))
            elif name == "list_scheduled_tasks":
                return self._list_scheduled_tasks()
            elif name == "cancel_task":
                return self._cancel_task(arguments.get("job_id"))
            elif name == "create_pdf":
                return self._create_pdf(arguments.get("filename"), arguments.get("content"))
            elif name == "create_docx":
                return self._create_docx(arguments.get("filename"), arguments.get("content"))
            elif name == "create_ppt":
                return self._create_ppt(arguments.get("filename"), arguments.get("title"), arguments.get("slides"))
            elif name == "create_excel":
                return self._create_excel(arguments.get("filename"), arguments.get("data"))
            elif name == "generate_linkedin_post":
                return await self._generate_linkedin_post(arguments.get("topic"))
            elif name == "post_to_linkedin":
                return await self._post_to_linkedin(arguments.get("text"), arguments.get("image_filenames", []))
            # === BROWSER AUTOMATION TOOLS (BrowserOS Architecture) ===
            elif name == "open_browser":
                return await browser_automation.open_browser(arguments.get("url"))
            elif name == "take_snapshot":
                return await browser_automation.take_snapshot()
            elif name == "click":
                return await browser_automation.click(arguments.get("uid"), arguments.get("dbl_click", False))
            elif name == "hover":
                return await browser_automation.hover(arguments.get("uid"))
            elif name == "fill":
                return await browser_automation.fill(arguments.get("uid"), arguments.get("value"))
            elif name == "fill_form":
                return await browser_automation.fill_form(arguments.get("elements", []))
            elif name == "select_option":
                return await browser_automation.select_option(arguments.get("uid"), arguments.get("option_text"))
            elif name == "drag":
                return await browser_automation.drag(arguments.get("from_uid"), arguments.get("to_uid"))
            elif name == "upload_file":
                return await browser_automation.upload_file(arguments.get("uid"), arguments.get("file_path"))
            elif name == "navigate_page":
                return await browser_automation.navigate_page(arguments.get("url"))
            elif name == "navigate_history":
                return await browser_automation.navigate_history(arguments.get("direction"))
            elif name == "new_page":
                return await browser_automation.new_page(arguments.get("url"))
            elif name == "list_pages":
                return await browser_automation.list_pages()
            elif name == "select_page":
                return await browser_automation.select_page(arguments.get("page_idx"))
            elif name == "close_page":
                return await browser_automation.close_page(arguments.get("page_idx"))
            elif name == "close_browser":
                return await browser_automation.close_browser()
            elif name == "take_screenshot":
                return await browser_automation.take_screenshot(arguments.get("uid"), arguments.get("full_page", False))
            elif name == "extract_text":
                return await browser_automation.extract_text()
            elif name == "extract_structured_data":
                return await browser_automation.extract_structured_data(arguments.get("data_type", "auto"))
            elif name == "type_text":
                return await browser_automation.type_text(arguments.get("text"), arguments.get("uid"))
            elif name == "press_key":
                return await browser_automation.press_key(arguments.get("key"), arguments.get("modifiers"))
            elif name == "scroll_page":
                return await browser_automation.scroll_page(arguments.get("direction", "down"))
            elif name == "scroll_to_element":
                return await browser_automation.scroll_to_element(arguments.get("uid"))
            elif name == "wait_for":
                return await browser_automation.wait_for(arguments.get("text"), arguments.get("timeout", 5000))
            elif name == "wait_for_navigation":
                return await browser_automation.wait_for_navigation(arguments.get("timeout", 10000))
            elif name == "execute_javascript":
                return await browser_automation.execute_javascript(arguments.get("code"))
            elif name == "handle_dialog":
                return await browser_automation.handle_dialog(arguments.get("action", "accept"), arguments.get("prompt_text"))
            elif name == "switch_to_frame":
                return await browser_automation.switch_to_frame(arguments.get("uid"))
            elif name == "switch_to_main":
                return await browser_automation.switch_to_main()
            elif name == "get_page_info":
                return await browser_automation.get_page_info()
            elif name == "submit_form":
                return await browser_automation.submit_form()
            # Legacy tool name compatibility
            elif name == "get_page_elements":
                return await browser_automation.take_snapshot()
            elif name == "click_element":
                return await browser_automation.click(arguments.get("uid") or arguments.get("selector"))
            elif name == "fill_input":
                return await browser_automation.fill(arguments.get("uid") or arguments.get("selector"), arguments.get("value"))
            elif name == "hover_element":
                return await browser_automation.hover(arguments.get("uid") or arguments.get("selector"))
            elif name == "navigate_to":
                return await browser_automation.navigate_page(arguments.get("url"))
            elif name == "go_back":
                return await browser_automation.navigate_history('back')
            elif name == "go_forward":
                return await browser_automation.navigate_history('forward')
            elif name == "open_new_tab":
                return await browser_automation.new_page(arguments.get("url"))
            elif name == "switch_tab":
                return await browser_automation.select_page(arguments.get("index", 0))
            elif name == "close_tab":
                return await browser_automation.close_tab()
            elif name == "drag_and_drop":
                return await browser_automation.drag(arguments.get("from_uid") or arguments.get("source_selector"), arguments.get("to_uid") or arguments.get("target_selector"))
            elif name == "take_page_screenshot":
                return await browser_automation.take_screenshot()
            elif name == "wait_for_element":
                return await browser_automation.wait_for(arguments.get("text") or arguments.get("selector"), arguments.get("timeout", 5000))
            # === YOUTUBE TRANSCRIPT TOOLS ===
            elif name == "youtube_transcript_search":
                return self._youtube_transcript_search(arguments.get("video_url"), arguments.get("search_phrase"))
            elif name == "get_youtube_transcript":
                return self._get_youtube_transcript(arguments.get("video_url"), arguments.get("max_chars", 2000))
            else:
                return f"Error: Tool '{name}' not found."
        except Exception as e:
            return f"Action Failure: {str(e)}"

    # ... existing _real_search, _browse_url, _write_file ...



    # ... (existing init) ...

    def _get_setting(self, key: str) -> str:
        """Fetches a setting from the database, falling back to env vars."""
        try:
            db = SessionLocal()
            setting = db.query(models.SystemSetting).filter(models.SystemSetting.key == key).first()
            if setting and setting.value:
                return setting.value
        except Exception as e:
            print(f"DB Setting Error: {e}")
        finally:
            db.close()
        
        # Fallback to env
        return os.getenv(key)

    # Email draft storage (in-memory for simplicity)
    _email_draft = None

    def _draft_email(self, recipient: str, subject: str, body: str, attachments: List[str] = None) -> str:
        """Drafts an email and returns a preview for user approval."""
        user_name = self._get_setting("USER_NAME") or "User"
        
        # Store draft for confirmation
        MCPService._email_draft = {
            "recipient": recipient,
            "subject": subject,
            "body": body,
            "attachments": attachments or []
        }
        
        att_str = ", ".join(attachments) if attachments else "None"
        
        # Return formatted preview
        preview = f"""
📧 **Email Draft Preview**

**From:** {user_name}
**To:** {recipient}
**Subject:** {subject}
**Attachments:** {att_str}

**Message:**
{body}

---
Would you like me to send this email? You can:
- Say "yes" or "send it" to confirm
- Request changes like "change the subject to..." or "make it more formal"
- Say "cancel" to discard
"""
        return preview

    def _confirm_send_email(self, confirmed: bool) -> str:
        """Sends the drafted email after user confirmation."""
        if not confirmed:
            MCPService._email_draft = None
            return "Email draft discarded."
        
        if not MCPService._email_draft:
            return "No email draft found. Please draft an email first."
        
        draft = MCPService._email_draft
        smtp_email = self._get_setting("SMTP_EMAIL")
        smtp_password = self._get_setting("SMTP_PASSWORD")
        
        if not smtp_email or not smtp_password:
            return f"[SIMULATION] Email to {draft['recipient']} with subject '{draft['subject']}' logged (Credentials missing in Settings)."

        try:
            msg = EmailMessage()
            msg.set_content(draft['body'])
            msg["Subject"] = draft['subject']
            msg["From"] = smtp_email
            msg["To"] = draft['recipient']
            
            # Attach files
            for filename in draft.get("attachments", []):
                path = os.path.join(os.getcwd(), "agent_files", filename)
                if os.path.exists(path):
                    ctype, encoding = mimetypes.guess_type(path)
                    if ctype is None or encoding is not None:
                        # No guess could be made, or the file is encoded (compressed), so
                        # use a generic bag-of-bits type.
                        ctype = 'application/octet-stream'
                    
                    maintype, subtype = ctype.split('/', 1)
                    
                    with open(path, 'rb') as f:
                        file_data = f.read()
                        msg.add_attachment(file_data,
                                           maintype=maintype,
                                           subtype=subtype,
                                           filename=filename)
                else:
                    return f"Error: Attachment '{filename}' not found."

            # Standard Gmail SMTP port 465 (SSL) or 587 (TLS)
            with smtplib.SMTP_SSL("smtp.gmail.com", 465) as smtp:
                smtp.login(smtp_email, smtp_password)
                smtp.send_message(msg)
            
            MCPService._email_draft = None
            return f"✅ Email sent successfully to {draft['recipient']}!"
        except Exception as e:
            return f"Email Error: {str(e)}"

    def _schedule_task(self, task_description: str, interval_seconds: int) -> str:
        """Schedules a task via the SchedulerService."""
        try:
            from app.services.scheduler_service import SchedulerService
            service = SchedulerService()
            job_id = service.add_job(task_description, "interval", str(interval_seconds))
            if job_id:
                return f"Task '{task_description}' scheduled successfully (ID: {job_id}, Interval: {interval_seconds}s)."
            else:
                return "Failed to schedule task."
        except Exception as e:
            return f"Scheduling Error: {str(e)}"

    def _list_scheduled_tasks(self) -> str:
        """Lists active jobs using SchedulerService."""
        try:
            from app.services.scheduler_service import SchedulerService
            service = SchedulerService()
            jobs = service.list_jobs()
            
            if not jobs:
                return "No active scheduled tasks found."
            
            output = ["📅 **Active Scheduled Tasks:**"]
            for job in jobs:
                next_run = "Pending"
                if job.get("next_run"):
                    # Format ISO string nicely if possible, or leave as is
                    next_run = job["next_run"].replace("T", " ")[:16]
                
                output.append(f"- **ID:** `{job['id']}`\n  **Task:** {job['name']}\n  **Next Run:** {next_run}")
            
            return "\n\n".join(output)
        except Exception as e:
            return f"List Error: {str(e)}"

    def _cancel_task(self, job_id: str) -> str:
        """Cancels a scheduled task."""
        try:
            from app.services.scheduler_service import SchedulerService
            service = SchedulerService()
            success = service.remove_job(job_id)
            if success:
                return f"Task with ID '{job_id}' has been successfully cancelled."
            else:
                return f"Failed to cancel task '{job_id}'. It may not exist."
        except Exception as e:
            return f"Cancellation Error: {str(e)}"

    def _read_email(self, limit: int = 5) -> str:
        """Reads unread emails using IMAP."""
        smtp_email = self._get_setting("SMTP_EMAIL")
        smtp_password = self._get_setting("SMTP_PASSWORD")
        
        if not smtp_email or not smtp_password:
            return "[SIMULATION] Checked inbox - No credentials found. Please set SMTP_EMAIL and SMTP_PASSWORD in Settings."

        try:
            # Connect to Gmail IMAP
            mail = imaplib.IMAP4_SSL("imap.gmail.com")
            mail.login(smtp_email, smtp_password)
            mail.select("inbox")

            # Search for unread emails
            status, messages = mail.search(None, "UNSEEN")
            if status != "OK":
                return "Failed to search emails."
            
            email_ids = messages[0].split()
            if not email_ids:
                return "No unread emails found."
            
            # Fetch latest `limit` emails
            latest_email_ids = email_ids[-limit:]
            email_list = []

            for e_id in reversed(latest_email_ids):
                status, msg_data = mail.fetch(e_id, "(RFC822)")
                for response_part in msg_data:
                    if isinstance(response_part, tuple):
                        msg = email.message_from_bytes(response_part[1])
                        subject, encoding = email.header.decode_header(msg["Subject"])[0]
                        if isinstance(subject, bytes):
                            subject = subject.decode(encoding or "utf-8")
                        sender = msg.get("From")
                        email_list.append(f"- From: {sender}\n  Subject: {subject}")
            
            mail.close()
            mail.logout()
            
            return "📧 **Unread Emails:**\n\n" + "\n\n".join(email_list)

        except Exception as e:
            return f"IMAP Error: {str(e)}"

    async def _real_search(self, query: str) -> str:
        """Attempts to perform a real search via Tavily or Serper."""
        # TAVILY IS PRIORITIZED FOR RICH RESULTS
        if self.tavily_api_key:
            try:
                async with httpx.AsyncClient() as client:
                    response = await client.post(
                        "https://api.tavily.com/search",
                        json={"api_key": self.tavily_api_key.strip(), "query": query, "search_depth": "basic"},
                        timeout=10.0
                    )
                    if response.status_code == 200:
                        results = response.json()
                        snippets = [r.get("content", "") for r in results.get("results", [])[:3]]
                        return f"Search Results (via Tavily):\n" + "\n".join(snippets)
            except Exception as e:
                print(f"Tavily Error: {e}")

        if self.serper_api_key:
            try:
                async with httpx.AsyncClient() as client:
                    headers = {"X-API-KEY": self.serper_api_key, "Content-Type": "application/json"}
                    response = await client.post("https://google.serper.dev/search", json={"q": query}, headers=headers)
                    if response.status_code == 200:
                        results = response.json()
                        snippets = [res.get("snippet", "") for res in results.get("organic", [])[:3]]
                        return "Search Results (via Serper):\n" + "\n".join(snippets)
            except Exception as e:
                print(f"Serper Error: {e}")

        # FALLBACK: Simple Fallback or Error
        lower_q = query.lower()
        if "bitcoin" in lower_q:
             return "Real-time Search Update: Bitcoin is trading at $96,520 USD as of the latest market tick. (Live via Simulation Engine)"
        
        return f"System Info: Performing real-time intelligence gathering for '{query}'. No live results found (check API keys in .env). Please provide a valid TAVILY_API_KEY or SERPER_API_KEY to fetch real-world data."

    def _sync_browse_url(self, url: str) -> str:
        """Sync implementation of browse_url - runs in thread pool."""
        try:
            from playwright.sync_api import sync_playwright
            
            with sync_playwright() as p:
                # Launch Visible Browser
                browser = p.chromium.launch(headless=False)
                # Maximize or set reasonable size
                context = browser.new_context(viewport={"width": 1280, "height": 800})
                page = context.new_page()
                
                print(f"Browsing: {url}")
                # Increased timeout and use domcontentloaded for faster response
                page.goto(url, timeout=90000, wait_until='domcontentloaded')
                
                # Wait for some content (simulates reading)
                page.wait_for_timeout(3000) 
                
                # Extract text
                text = page.evaluate("document.body.innerText")
                title = page.title()
                
                browser.close()
                
                # Basic cleaning & strict truncation
                cleaned = re.sub(r'\s+', ' ', text).strip()
                return f"Browsed: {title} | {url}\n{cleaned[:2000]}"
        except Exception as e:
            import traceback
            error_details = traceback.format_exc()
            print(f"Browse Error: {error_details}")
            return f"Playwright Browse Error: {str(e)}"
    
    async def _browse_url(self, url: str) -> str:
        """Browse URL using Playwright (Visible Mode) for handling dynamic content."""
        import concurrent.futures
        loop = asyncio.get_event_loop()
        with concurrent.futures.ThreadPoolExecutor() as executor:
            return await loop.run_in_executor(executor, self._sync_browse_url, url)

    def _sync_take_screenshot(self, url: str, filename: str) -> str:
        """Sync implementation of take_screenshot - runs in thread pool."""
        try:
            from playwright.sync_api import sync_playwright
            
            path = os.path.join(os.getcwd(), "agent_files", filename)
            os.makedirs(os.path.dirname(path), exist_ok=True)
            
            with sync_playwright() as p:
                browser = p.chromium.launch(headless=False)
                page = browser.new_page()
                page.goto(url, timeout=90000, wait_until='domcontentloaded')
                page.wait_for_timeout(2000)
                
                page.screenshot(path=path)
                browser.close()
                
            return f"Screenshot saved to '{filename}'."
        except Exception as e:
            import traceback
            error_details = traceback.format_exc()
            print(f"Screenshot Error: {error_details}")
            return f"Screenshot Error: {str(e)}"
    
    async def _take_screenshot(self, url: str, filename: str) -> str:
        """Takes a screenshot of the given URL."""
        import concurrent.futures
        loop = asyncio.get_event_loop()
        with concurrent.futures.ThreadPoolExecutor() as executor:
            return await loop.run_in_executor(executor, self._sync_take_screenshot, url, filename)

    def _write_file(self, filename: str, content: str) -> str:
        # Save to agent_files
        path = os.path.join(os.getcwd(), "agent_files", filename)
        os.makedirs(os.path.dirname(path), exist_ok=True)
        with open(path, "w") as f:
            f.write(content)
        return f"File '{filename}' written successfully."



    def _analyze_data(self, filename: str, query: str) -> str:
        """Reads CSV/Excel and returns a summary or analysis."""
        try:
            # DEBUG
            print(f"DEBUG: Analyze requested for '{filename}'")
            print(f"DEBUG: CWD is {os.getcwd()}")
            
            path = os.path.join(os.getcwd(), "agent_files", filename)
            print(f"DEBUG: Full path check: {path}")
            
            if not os.path.exists(path):
                return f"Error: File '{filename}' not found at {path}."
            
            # Detect type
            if filename.endswith(".csv"):
                df = pd.read_csv(path)
            elif filename.endswith(".xlsx"):
                df = pd.read_excel(path)
            else:
                return "Error: Unsupported file format. Use CSV or Excel."

            # Basic analysis
            info_buf = StringIO()
            df.info(buf=info_buf)
            info_str = info_buf.getvalue()
            
            desc = df.describe().to_markdown()
            head = df.head(5).to_markdown()
            
            result = (
                f"Data Analysis Results for {filename} (Truncated for efficiency):\n"
                f"### Stats Summary\n{desc[:1000]}\n"
                f"### Preview (5 Rows)\n{head}\n"
            )
            print(f"DEBUG: Returning result len={len(result)}")
            return result
        except Exception as e:
            print(f"DEBUG: Exception in analyze_data: {e}")
            return f"Data Analysis Error: {str(e)}"

    def _read_pdf(self, filename: str) -> str:
        """Extracts text from a PDF file."""
        try:
            path = os.path.join(os.getcwd(), "agent_files", filename)
            if not os.path.exists(path):
                return f"Error: File '{filename}' not found."
            
            reader = pypdf.PdfReader(path)
            text = []
            for page in reader.pages:
                text.append(page.extract_text())
            
            full_text = "\n".join(text)
            # Cap at 3000 chars (~750 tokens)
            if len(full_text) > 3000:
                return f"PDF: {filename}\n{full_text[:3000]}\n... [truncated]"
            return f"PDF: {filename}\n{full_text}"
        except Exception as e:
            return f"PDF Read Error: {str(e)}"

    def _create_pdf(self, filename: str, content: str) -> str:
        """Creates a PDF file with Markdown formatting support."""
        try:
            path = os.path.join(os.getcwd(), "agent_files", filename)
            os.makedirs(os.path.dirname(path), exist_ok=True)
            
            doc = SimpleDocTemplate(path, pagesize=letter)
            styles = getSampleStyleSheet()
            # Custom Styles
            if 'Code' not in styles:
                styles.add(ParagraphStyle(name='Code', parent=styles['Normal'], fontName='Courier', fontSize=9, leading=11, backColor=colors.lightgrey, borderPadding=5))
            if 'Quote' not in styles:
                styles.add(ParagraphStyle(name='Quote', parent=styles['Normal'], leftIndent=20, textColor=colors.darkgrey, borderPadding=2))
            if 'TableHeader' not in styles:
                 styles.add(ParagraphStyle(name='TableHeader', parent=styles['Normal'], textColor=colors.whitesmoke, fontName='Helvetica-Bold', alignment=TA_CENTER))

            story = []

            # Helper to build table
            def build_table(data):
                if not data: return None
                try:
                    num_cols = len(data[0])
                    avail_width = 460
                    col_width = avail_width / num_cols if num_cols else 0
                    
                    tbl_data = []
                    # Header
                    tbl_data.append([Paragraph(str(cell), styles['TableHeader']) for cell in data[0]])
                    # Body
                    for row in data[1:]:
                        # Ensure row length matches header
                        row_data = []
                        for i in range(num_cols):
                            cell = row[i] if i < len(row) else ""
                            row_data.append(Paragraph(str(cell), styles['Normal']))
                        tbl_data.append(row_data)

                    t = Table(tbl_data, colWidths=[col_width]*num_cols, repeatRows=1)
                    t.setStyle(TableStyle([
                        ('VALIGN', (0,0), (-1,-1), 'TOP'),
                        ('BACKGROUND', (0, 0), (-1, 0), colors.Color(0.2, 0.4, 0.6)),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black)
                    ]))
                    return t
                except Exception as e:
                    return Paragraph(f"Error building table: {e}", styles['Normal'])

            # Helper to format inline markdown (Bold, Italic, Code)
            def format_inline(text):
                # 1. Double backticks for literal backticks or generic code: ``text`` -> <font name="Courier" backColor="lightgrey">text</font>
                text = re.sub(r'``(.*?)``', r'<font name="Courier" backColor="lightgrey"> \1 </font>', text)
                # 2. Single backticks: `text` -> <font name="Courier" backColor="lightgrey">text</font>
                text = re.sub(r'`(.*?)`', r'<font name="Courier" backColor="lightgrey"> \1 </font>', text)
                # 3. Bold: **text** -> <b>text</b>
                text = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', text)
                # 4. Italic: *text* -> <i>text</i>
                text = re.sub(r'(?<!\*)\*(.*?)\*(?!\*)', r'<i>\1</i>', text)
                return text

            # 1. Split lines
            lines = content.split('\n')
            current_table_data = []
            in_code_block = False
            code_buffer = []

            for line in lines:
                raw_line = line # Keep raw for code blocks
                line = line.strip()

                # --- Code Block Handling ---
                if line.startswith('```'):
                    if in_code_block:
                        # End of block
                        full_code = "\n".join(code_buffer)
                        story.append(Preformatted(full_code, styles["Code"]))
                        story.append(Spacer(1, 12))
                        code_buffer = []
                        in_code_block = False
                    else:
                        # Start of block
                        # Flush any pending table first
                        if current_table_data:
                            t = build_table(current_table_data)
                            if t:
                                story.append(t)
                                story.append(Spacer(1, 12))
                            current_table_data = []
                        in_code_block = True
                    continue 

                if in_code_block:
                    code_buffer.append(raw_line)
                    continue

                if not line:
                    # Flush table if exists
                    if current_table_data:
                        t = build_table(current_table_data)
                        if t:
                            story.append(t)
                            story.append(Spacer(1, 12))
                        current_table_data = []
                    continue
                
                # --- Table Handling ---
                if line.startswith('|') and line.endswith('|'):
                    row = [cell.strip() for cell in line[1:-1].split('|')]
                    if not all(c.replace('-', '') == '' for c in row): # Ignore divider |---|
                         # Format cells!
                        formatted_row = [format_inline(c) for c in row]
                        current_table_data.append(formatted_row)
                    continue
                
                # If we hit a non-table line, flush table
                if current_table_data:
                    t = build_table(current_table_data)
                    if t:
                        story.append(t)
                        story.append(Spacer(1, 12))
                    current_table_data = []

                # --- Element Parsing ---
                
                # Apply inline formatting ONLY for non-code text
                # We do this here so it doesn't affect the table/code logic checks above
                
                # Headers
                if line.startswith('# '):
                    p = Paragraph(format_inline(line[2:]), styles["Heading1"])
                elif line.startswith('## '):
                    p = Paragraph(format_inline(line[3:]), styles["Heading2"])
                elif line.startswith('### '):
                    p = Paragraph(format_inline(line[4:]), styles["Heading3"])
                
                # List Items (Unordered)
                elif line.startswith('- ') or line.startswith('* '):
                    p = Paragraph(f"•  {format_inline(line[2:])}", styles["Normal"])
                
                # List Items (Ordered)
                elif re.match(r'^\d+\.\s', line):
                    p = Paragraph(format_inline(line), styles["Normal"])

                # Blockquotes
                elif line.startswith('> '):
                    p = Paragraph(format_inline(line[2:]), styles["Quote"])
                
                # Normal Text
                else:
                    p = Paragraph(format_inline(line), styles["Normal"])
                
                story.append(p)
                story.append(Spacer(1, 8))
            
            # Flush final table/code
            if current_table_data:
                t = build_table(current_table_data)
                if t:
                    story.append(t)
            
            if code_buffer: # Unclosed code block
                 full_code = "\n".join(code_buffer)
                 story.append(Preformatted(full_code, styles["Code"]))

            doc.build(story)
            return f"PDF '{filename}' created successfully with rich Markdown (Tables, Code, Lists)."
        except Exception as e:
            return f"Create PDF Error: {str(e)}"

    def _create_docx(self, filename: str, content: str) -> str:
        """Creates a Word document with Markdown formatting support."""
        try:
            path = os.path.join(os.getcwd(), "agent_files", filename)
            os.makedirs(os.path.dirname(path), exist_ok=True)
            
            doc = Document()
            
            lines = content.split('\n')
            current_table_data = []
            
            for line in lines:
                line = line.strip()
                if not line:
                    # Flush table
                    if current_table_data:
                        rows = len(current_table_data)
                        cols = len(current_table_data[0])
                        table = doc.add_table(rows=rows, cols=cols)
                        table.style = 'Table Grid'
                        for i, row_data in enumerate(current_table_data):
                            for j, cell_text in enumerate(row_data):
                                table.cell(i, j).text = cell_text
                        doc.add_paragraph() # Spacer
                        current_table_data = []
                    continue
                
                # Check for table
                if line.startswith('|') and line.endswith('|'):
                    row = [cell.strip() for cell in line[1:-1].split('|')]
                    # Ignore separator rows
                    if not all(c.replace('-', '') == '' for c in row):
                        current_table_data.append(row)
                    continue

                if current_table_data:
                    rows = len(current_table_data)
                    cols = len(current_table_data[0])
                    table = doc.add_table(rows=rows, cols=cols)
                    table.style = 'Table Grid'
                    for i, row_data in enumerate(current_table_data):
                        for j, cell_text in enumerate(row_data):
                            table.cell(i, j).text = cell_text
                    doc.add_paragraph() # Spacer
                    current_table_data = []

                # Headers
                if line.startswith('# '):
                    p = doc.add_heading(line[2:], level=1)
                elif line.startswith('## '):
                    p = doc.add_heading(line[3:], level=2)
                else:
                    p = doc.add_paragraph()
                    # Parse bold
                    # Example: "This is **bold** text" -> ["This is ", "bold", " text"]
                    parts = re.split(r'(\*\*.*?\*\*)', line)
                    for part in parts:
                        if part.startswith('**') and part.endswith('**'):
                            run = p.add_run(part[2:-2])
                            run.bold = True
                        else:
                            p.add_run(part)

            # Flush final table
            if current_table_data:
                  rows = len(current_table_data)
                  cols = len(current_table_data[0])
                  table = doc.add_table(rows=rows, cols=cols)
                  table.style = 'Table Grid'
                  for i, row_data in enumerate(current_table_data):
                      for j, cell_text in enumerate(row_data):
                          table.cell(i, j).text = cell_text

            doc.save(path)
            return f"Word document '{filename}' created successfully with structured tables."
        except Exception as e:
            return f"Create Docx Error: {str(e)}"

    def _create_ppt(self, filename: str, title: str, slides: List[Dict[str, Any]]) -> str:
        """Creates a PowerPoint presentation."""
        try:
            path = os.path.join(os.getcwd(), "agent_files", filename)
            os.makedirs(os.path.dirname(path), exist_ok=True)
            
            prs = Presentation()
            
            # Title Slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            title_shape.text = title
            subtitle_shape.text = "Generated by Edith"
            
            # Content Slides
            bullet_slide_layout = prs.slide_layouts[1]
            for slide_data in slides:
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                title_shape.text = slide_data.get("title", "Slide")
                tf = body_shape.text_frame
                
                content_list = slide_data.get("content", [])
                if content_list:
                    tf.text = content_list[0]
                    for item in content_list[1:]:
                        p = tf.add_paragraph()
                        p.text = item

            prs.save(path)
            return f"PowerPoint '{filename}' created successfully."
        except Exception as e:
            return f"Create PPT Error: {str(e)}"

    def _create_excel(self, filename: str, data: List[Dict[str, Any]]) -> str:
        """Creates an Excel spreadsheet with styled headers."""
        try:
            path = os.path.join(os.getcwd(), "agent_files", filename)
            os.makedirs(os.path.dirname(path), exist_ok=True)
            
            if not data:
                return "Error: No data provided for Excel file."
                
            df = pd.DataFrame(data)
            
            # Use ExcelWriter with XlsxWriter engine for styling
            with pd.ExcelWriter(path, engine='xlsxwriter') as writer:
                df.to_excel(writer, index=False, sheet_name='Sheet1')
                workbook  = writer.book
                worksheet = writer.sheets['Sheet1']
                
                # Add a header format
                header_format = workbook.add_format({
                    'bold': True,
                    'text_wrap': True,
                    'valign': 'top',
                    'fg_color': '#D7E4BC',
                    'border': 1
                })
                
                # Write the column headers with the defined format
                for col_num, value in enumerate(df.columns.values):
                    worksheet.write(0, col_num, value, header_format)
                    pass # Values are valid

            return f"Excel file '{filename}' created successfully with styled headers."
        except Exception as e:
            return f"Create Excel Error: {str(e)}"

    async def _generate_linkedin_post(self, topic: str) -> str:
        """Generate a professional LinkedIn post using LLM"""
        try:
            from app.services.llm_service import llm_service
            
            prompt = f"""Generate a professional LinkedIn post about the following topic. 
            Make it engaging, authentic, and suitable for a professional audience.
            Use appropriate emojis sparingly, include relevant hashtags, and keep it under 1500 characters.
            
            Topic: {topic}
            
            Generate ONLY the post text, nothing else."""
            
            response = await llm_service.get_response(prompt)
            return f"""✅ LinkedIn Post Generated Successfully!

{response}

To publish this post, use the 'post_to_linkedin' tool with the text above."""
        except Exception as e:
            return f"Post Generation Error: {str(e)}"

    async def _post_to_linkedin(self, text: str, image_filenames: List[str] = None) -> str:
        """Post content to LinkedIn with optional images"""
        try:
            from app.services.linkedin_service import linkedin_service
            
            # Check authentication
            if not linkedin_service.is_authenticated():
                auth_url = linkedin_service.get_authorization_url()
                return f"""LinkedIn Authentication Required!
                
Please visit this URL to authorize EDITH:
{auth_url}

After authorization, try posting again."""
            
            # Upload images if provided
            image_urns = []
            if image_filenames:
                for filename in image_filenames:
                    # Check if it's an absolute path (from chat upload) or relative (in agent_files)
                    if os.path.isabs(filename) and os.path.exists(filename):
                        image_path = filename
                    else:
                        image_path = os.path.join(os.getcwd(), "agent_files", filename)
                        
                    if not os.path.exists(image_path):
                        return f"Error: Image '{filename}' not found (checked absolute path and agent_files)."
                    
                    urn = await linkedin_service.upload_image(image_path)
                    if urn:
                        image_urns.append(urn)
                    else:
                        return f"Error: Failed to upload image '{filename}'."
            
            # Create post
            result = await linkedin_service.create_post(text, image_urns if image_urns else None)
            
            if result.get("success"):
                return f"✅ Successfully posted to LinkedIn!\n\nPost ID: {result.get('post_id')}"
            else:
                return f"❌ Failed to post to LinkedIn: {result.get('error')}"
                
        except Exception as e:
            return f"LinkedIn Post Error: {str(e)}"

    # === YOUTUBE TRANSCRIPT TOOLS ===

    def _extract_video_id(self, url: str) -> str:
        """Extracts YouTube video ID from various URL formats."""
        patterns = [
            r'(?:v=|/v/|youtu\.be/)([a-zA-Z0-9_-]{11})',
            r'(?:embed/)([a-zA-Z0-9_-]{11})',
            r'^([a-zA-Z0-9_-]{11})$',  # Just the ID
        ]
        for pattern in patterns:
            match = stdlib_re.search(pattern, url)
            if match:
                return match.group(1)
        return None

    def _fetch_transcript(self, video_id: str):
        """Fetches transcript using v1.2.4 API. Returns list of FetchedTranscriptSnippet."""
        from youtube_transcript_api import YouTubeTranscriptApi
        ytt = YouTubeTranscriptApi()
        tlist = ytt.list(video_id)
        # Try manual captions first, then auto-generated
        try:
            t = tlist.find_transcript(['en'])
        except:
            t = tlist.find_generated_transcript(['en'])
        return t.fetch()

    def _youtube_transcript_search(self, video_url: str, search_phrase: str) -> str:
        """Searches a YouTube video transcript for a phrase. Returns timestamp + context."""
        try:
            video_id = self._extract_video_id(video_url)
            if not video_id:
                return f"Could not extract video ID from: {video_url}"

            try:
                snippets = self._fetch_transcript(video_id)
            except Exception as e:
                return f"No transcript available for this video. Error: {str(e)}"

            # Build full text with timestamps for searching
            search_lower = search_phrase.lower()
            
            # Search in sliding window of combined entries
            matches = []
            for i in range(len(snippets)):
                window_texts = []
                for j in range(i, min(i + 5, len(snippets))):
                    window_texts.append(snippets[j].text)
                
                combined = ' '.join(window_texts).lower()
                combined = stdlib_re.sub(r'\s+', ' ', combined)
                
                if search_lower in combined:
                    seconds = int(snippets[i].start)
                    minutes = seconds // 60
                    secs = seconds % 60
                    context = ' '.join(window_texts)
                    
                    matches.append({
                        'timestamp': f"{minutes}:{secs:02d}",
                        'seconds': seconds,
                        'context': context[:200]
                    })
                    break

            if not matches:
                # Fallback: fuzzy search using word overlap
                search_words = set(search_lower.split())
                best_score = 0
                best_match = None
                
                for i in range(len(snippets)):
                    window_texts = []
                    for j in range(i, min(i + 5, len(snippets))):
                        window_texts.append(snippets[j].text)
                    combined = ' '.join(window_texts).lower()
                    combined_words = set(combined.split())
                    
                    overlap = len(search_words & combined_words) / max(len(search_words), 1)
                    if overlap > best_score and overlap > 0.5:
                        best_score = overlap
                        seconds = int(snippets[i].start)
                        best_match = {
                            'timestamp': f"{seconds // 60}:{seconds % 60:02d}",
                            'seconds': seconds,
                            'context': ' '.join(window_texts)[:200],
                            'score': f"{overlap:.0%}"
                        }
                
                if best_match:
                    return (
                        f"Closest match ({best_match['score']} confidence) at {best_match['timestamp']} "
                        f"(t={best_match['seconds']}s)\n"
                        f"Context: \"{best_match['context']}\"\n"
                        f"Play from here: {video_url}&t={best_match['seconds']}"
                    )
                return f"Phrase not found in transcript. Try different wording or a shorter phrase."

            m = matches[0]
            return (
                f"Found at {m['timestamp']} (t={m['seconds']}s)\n"
                f"Context: \"{m['context']}\"\n"
                f"Play from here: {video_url}&t={m['seconds']}"
            )

        except Exception as e:
            return f"YouTube transcript error: {str(e)}"

    def _get_youtube_transcript(self, video_url: str, max_chars: int = 2000) -> str:
        """Fetches YouTube transcript with a character limit to save tokens."""
        try:
            video_id = self._extract_video_id(video_url)
            if not video_id:
                return f"Could not extract video ID from: {video_url}"

            try:
                snippets = self._fetch_transcript(video_id)
            except Exception as e:
                return f"No transcript available: {str(e)}"

            # Format with timestamps, respecting max_chars
            output = ""
            for s in snippets:
                seconds = int(s.start)
                timestamp = f"{seconds // 60}:{seconds % 60:02d}"
                line = f"[{timestamp}] {s.text}\n"
                
                if len(output) + len(line) > max_chars:
                    output += f"\n... (truncated at {max_chars} chars. Use youtube_transcript_search for specific phrases)"
                    break
                output += line

            return output

        except Exception as e:
            return f"YouTube transcript error: {str(e)}"

mcp_service = MCPService()

